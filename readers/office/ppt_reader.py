"""
PowerPoint 文档读取器
"""
from pathlib import Path

from core.file_reader import BaseFileReader, FileContent
from config import CHUNK_SIZE, CHUNK_OVERLAP


class PptReader(BaseFileReader):
    """PowerPoint 读取器 (.pptx)"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = [".pptx", ".ppt"]
    
    def read(self, file_path: Path) -> FileContent:
        """读取 PowerPoint 文档"""
        file_path = Path(file_path)
        
        prs = None
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [f"=== 幻灯片 {i} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                if len(slide_texts) > 1:
                    slides_content.append("\n".join(slide_texts))
            
            content = "\n\n".join(slides_content)
            title = prs.core_properties.title or file_path.stem
            
        except Exception as e:
            content = f"[Error reading PowerPoint: {e}]"
            title = file_path.stem
        
        # 提取元数据
        metadata = self.extract_metadata_from_path(file_path)
        metadata["slide_count"] = len(prs.slides) if prs else 0
        
        # 分块
        chunks = self.chunk_content(content, CHUNK_SIZE, CHUNK_OVERLAP)
        
        return FileContent(
            source=str(file_path.absolute()),
            content=content,
            file_type="powerpoint",
            title=title,
            metadata=metadata,
            chunks=chunks,
            file_hash=self.calculate_hash(content)
        )
