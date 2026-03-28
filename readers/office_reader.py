"""
Office 文档读取器 - Word, PowerPoint, Excel
"""
from pathlib import Path
from typing import List

from core.file_reader import BaseFileReader, FileContent
from config import CHUNK_SIZE, CHUNK_OVERLAP


class WordReader(BaseFileReader):
    """Word 文档读取器 (.docx)"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = [".docx", ".doc"]
    
    def read(self, file_path: Path) -> FileContent:
        """读取 Word 文档"""
        file_path = Path(file_path)
        
        try:
            from docx import Document
            
            doc = Document(file_path)
            
            # 提取文本
            paragraphs = []
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        paragraphs.append(" | ".join(row_text))
            
            content = "\n\n".join(paragraphs)
            
            # 提取标题（使用第一个段落或文件名）
            title = doc.paragraphs[0].text if doc.paragraphs and doc.paragraphs[0].text else file_path.stem
            
        except Exception as e:
            # 如果 python-docx 失败，尝试作为纯文本读取
            content = f"[Error reading Word document: {e}]"
            title = file_path.stem
        
        # 提取元数据
        metadata = self.extract_metadata_from_path(file_path)
        
        # 分块
        chunks = self.chunk_content(content, CHUNK_SIZE, CHUNK_OVERLAP)
        
        return FileContent(
            source=str(file_path.absolute()),
            content=content,
            file_type="word",
            title=title[:100],  # 限制标题长度
            metadata=metadata,
            chunks=chunks,
            file_hash=self.calculate_hash(content)
        )


class PowerPointReader(BaseFileReader):
    """PowerPoint 读取器 (.pptx)"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = [".pptx", ".ppt"]
    
    def read(self, file_path: Path) -> FileContent:
        """读取 PowerPoint 文档"""
        file_path = Path(file_path)
        
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            
            slides_content = []
            
            for i, slide in enumerate(prs.slides, 1):
                slide_texts = [f"=== 幻灯片 {i} ==="]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_texts.append(shape.text)
                
                if len(slide_texts) > 1:
                    slides_content.append("\n".join(slide_texts))
            
            content = "\n\n".join(slides_content)
            title = prs.core_properties.title or file_path.stem
            
        except Exception as e:
            content = f"[Error reading PowerPoint: {e}]"
            title = file_path.stem
        
        # 提取元数据
        metadata = self.extract_metadata_from_path(file_path)
        metadata["slide_count"] = len(prs.slides) if 'prs' in dir() else 0
        
        # 分块
        chunks = self.chunk_content(content, CHUNK_SIZE, CHUNK_OVERLAP)
        
        return FileContent(
            source=str(file_path.absolute()),
            content=content,
            file_type="powerpoint",
            title=title,
            metadata=metadata,
            chunks=chunks,
            file_hash=self.calculate_hash(content)
        )


class ExcelReader(BaseFileReader):
    """Excel 读取器 (.xlsx, .xls)"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = [".xlsx", ".xls"]
    
    def read(self, file_path: Path) -> FileContent:
        """读取 Excel 文档"""
        file_path = Path(file_path)
        
        try:
            from openpyxl import load_workbook
            
            wb = load_workbook(file_path, data_only=True)
            
            all_sheets = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                sheet_rows = [f"=== 工作表: {sheet_name} ==="]
                
                for row in sheet.iter_rows():
                    row_values = []
                    for cell in row:
                        if cell.value is not None:
                            row_values.append(str(cell.value))
                    
                    if row_values:
                        sheet_rows.append(" | ".join(row_values))
                
                all_sheets.append("\n".join(sheet_rows))
            
            content = "\n\n".join(all_sheets)
            title = file_path.stem
            
        except Exception as e:
            content = f"[Error reading Excel: {e}]"
            title = file_path.stem
        
        # 提取元数据
        metadata = self.extract_metadata_from_path(file_path)
        metadata["sheet_count"] = len(wb.sheetnames) if 'wb' in dir() else 0
        
        # 分块
        chunks = self.chunk_content(content, CHUNK_SIZE, CHUNK_OVERLAP)
        
        return FileContent(
            source=str(file_path.absolute()),
            content=content,
            file_type="excel",
            title=title,
            metadata=metadata,
            chunks=chunks,
            file_hash=self.calculate_hash(content)
        )
