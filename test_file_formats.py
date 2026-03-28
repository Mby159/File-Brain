"""
测试各种文件格式的读取
"""
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

CHECK = "[OK]"
CROSS = "[X]"
SKIP = "[SKIP]"


def _test_format(reader_class, file_path, expected_type):
    """测试单个文件格式"""
    try:
        reader = reader_class()
        content = reader.read(file_path)
        
        if content.file_type == expected_type:
            print(f"{CHECK} {expected_type}: 读取成功")
            print(f"    标题: {content.title[:50]}...")
            print(f"    内容长度: {len(content.content)} 字符")
            return True
        else:
            print(f"{CROSS} {expected_type}: 类型不匹配 (期望 {expected_type}, 实际 {content.file_type})")
            return False
    except Exception as e:
        print(f"{CROSS} {expected_type}: 读取失败 - {e}")
        return False


def test_all_formats():
    """测试所有支持的文件格式"""
    print("=" * 60)
    print("测试文件格式读取")
    print("=" * 60)
    
    test_dir = Path(__file__).parent / "test_files"
    test_dir.mkdir(exist_ok=True)
    
    results = {}
    
    # 1. 测试纯文本
    print("\n【文本文件】")
    from readers.text_reader import TextReader, MarkdownReader
    
    txt_file = test_dir / "test.txt"
    txt_file.write_text("这是纯文本文件内容。\n用于测试 File Brain。", encoding='utf-8')
    results['txt'] = _test_format(TextReader, txt_file, "text")
    
    # 2. 测试 Markdown
    md_file = test_dir / "test.md"
    md_file.write_text("""---
title: Markdown测试
tags: [test]
---

# 标题

这是 Markdown 内容。
""", encoding='utf-8')
    results['md'] = _test_format(MarkdownReader, md_file, "markdown")
    
    # 3. 测试 Word (需要 python-docx)
    print("\n【Office 文档】")
    try:
        from readers.office_reader import WordReader
        from docx import Document
        
        docx_file = test_dir / "test.docx"
        doc = Document()
        doc.add_heading('Word 测试文档', 0)
        doc.add_paragraph('这是 Word 文档的内容。')
        doc.add_paragraph('用于测试 File Brain 的 Word 读取功能。')
        doc.save(docx_file)
        
        results['docx'] = _test_format(WordReader, docx_file, "word")
    except ImportError:
        print(f"{SKIP} docx: python-docx 未安装")
        results['docx'] = None
    except Exception as e:
        print(f"{CROSS} docx: {e}")
        results['docx'] = False
    
    # 4. 测试 PowerPoint
    try:
        from readers.office_reader import PowerPointReader
        from pptx import Presentation
        
        pptx_file = test_dir / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "PPT 测试"
        slide.placeholders[1].text = "这是 PowerPoint 内容"
        prs.save(pptx_file)
        
        results['pptx'] = _test_format(PowerPointReader, pptx_file, "powerpoint")
    except ImportError:
        print(f"{SKIP} pptx: python-pptx 未安装")
        results['pptx'] = None
    except Exception as e:
        print(f"{CROSS} pptx: {e}")
        results['pptx'] = False
    
    # 5. 测试 Excel
    try:
        from readers.office_reader import ExcelReader
        from openpyxl import Workbook
        
        xlsx_file = test_dir / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws.title = "测试表"
        ws['A1'] = "标题"
        ws['B1'] = "内容"
        ws['A2'] = "测试1"
        ws['B2'] = "Excel 内容"
        wb.save(xlsx_file)
        
        results['xlsx'] = _test_format(ExcelReader, xlsx_file, "excel")
    except ImportError:
        print(f"{SKIP} xlsx: openpyxl 未安装")
        results['xlsx'] = None
    except Exception as e:
        print(f"{CROSS} xlsx: {e}")
        results['xlsx'] = False
    
    # 6. 测试 PDF
    print("\n【PDF 文件】")
    try:
        from readers.others.pdf_reader import PdfReader
        
        # 创建一个简单的 PDF 文件用于测试
        from fpdf import FPDF
        
        pdf_file = test_dir / "test.pdf"
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt="PDF Test Document", ln=1, align="C")
        pdf.cell(200, 10, txt="This is PDF content.", ln=2)
        pdf.cell(200, 10, txt="For testing File Brain PDF reading functionality.", ln=3)
        pdf.output(str(pdf_file))
        
        results['pdf'] = _test_format(PdfReader, pdf_file, "pdf")
    except ImportError as e:
        print(f"{SKIP} pdf: {e}")
        results['pdf'] = None
    except Exception as e:
        print(f"{CROSS} pdf: {e}")
        results['pdf'] = False
    
    # 7. 测试 HTML
    print("\n【HTML 文件】")
    try:
        from readers.others.html_reader import HtmlReader
        
        html_file = test_dir / "test.html"
        html_file.write_text("""<!DOCTYPE html>
<html>
<head>
    <title>HTML 测试</title>
    <meta name="description" content="测试 HTML 读取">
</head>
<body>
    <h1>HTML 测试文档</h1>
    <p>这是 HTML 内容。</p>
    <article>
        <h2>文章标题</h2>
        <p>文章内容...</p>
    </article>
</body>
</html>""", encoding='utf-8')
        
        results['html'] = _test_format(HtmlReader, html_file, "html")
    except ImportError as e:
        print(f"{SKIP} html: {e}")
        results['html'] = None
    except Exception as e:
        print(f"{CROSS} html: {e}")
        results['html'] = False
    
    # 8. 测试思维导图
    print("\n【思维导图】")
    try:
        from readers.others.mindmap_reader import MindmapReader
        
        # 创建一个简单的 FreeMind 文件用于测试
        mm_file = test_dir / "test.mm"
        mm_content = """<?xml version="1.0" encoding="UTF-8"?>
<map version="1.0.1">
  <node TEXT="思维导图测试">
    <node TEXT="主题 1"/>
    <node TEXT="主题 2"/>
  </node>
</map>"""
        mm_file.write_text(mm_content, encoding='utf-8')
        
        results['mindmap'] = _test_format(MindmapReader, mm_file, "mindmap")
    except ImportError as e:
        print(f"{SKIP} xmind/mm: {e}")
        results['mindmap'] = None
    except Exception as e:
        print(f"{CROSS} mindmap: {e}")
        results['mindmap'] = False
    
    # 统计结果
    print("\n" + "=" * 60)
    print("测试结果统计")
    print("=" * 60)
    
    passed = sum(1 for v in results.values() if v is True)
    failed = sum(1 for v in results.values() if v is False)
    skipped = sum(1 for v in results.values() if v is None)
    total = len(results)
    
    print(f"通过: {passed}/{total}")
    print(f"失败: {failed}/{total}")
    print(f"跳过: {skipped}/{total}")
    
    print("\n详细结果:")
    for fmt, result in results.items():
        status = CHECK if result is True else (CROSS if result is False else SKIP)
        print(f"  {status} {fmt}")
    
    return failed == 0


if __name__ == "__main__":
    success = test_all_formats()
    sys.exit(0 if success else 1)
